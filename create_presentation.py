import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # 1. Initialize presentation
    prs = Presentation()
    
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 2. Define Theme Colors (Matching Samarkand Night Design System)
    C_BG_DARK = RGBColor(10, 15, 29)       # #0A0F1D (Dark Sapphire)
    C_BG_CARD = RGBColor(18, 26, 47)       # #121A2F (Glass Card BG)
    C_GOLD = RGBColor(212, 175, 55)        # #D4AF37 (Regal Gold)
    C_TURQUOISE = RGBColor(0, 155, 158)    # #009B9E (Deep Turquoise)
    C_WHITE = RGBColor(255, 255, 255)      # White
    C_SLATE_400 = RGBColor(148, 163, 184)  # #94A3B8 (Secondary Text)
    C_SLATE_600 = RGBColor(71, 85, 105)    # #475569 (Muted Text)

    # 3. Helper to style slide background
    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = C_BG_DARK
        
        # Add Uzbek Mandala Ornament peeking in the top-right corner
        try:
            slide.shapes.add_picture("uzbek_mandala.jpg", Inches(11.0), Inches(-1.0), Inches(3.5), Inches(3.5))
        except Exception as e:
            pass
        
        # Add top decorative golden line
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = C_GOLD
        top_line.line.fill.background()
        
        # Add bottom brand branding
        brand_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(10.5), Inches(0.3))
        tf = brand_box.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "UZBEK-TRIP (SAMARQAND CRAFTOUR)  |  5-DAQIQALIK INVESTOR PITCH DECK"
        p.font.name = "Segoe UI"
        p.font.size = Pt(9)
        p.font.color.rgb = C_SLATE_600
        p.font.bold = True
        
        # Add small logo in the bottom right corner
        try:
            slide.shapes.add_picture("logo.jpg", Inches(12.5), Inches(6.9), Inches(0.5), Inches(0.5))
        except Exception as e:
            pass

    # 4. Helper to create standard titles
    def add_slide_title(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = "Georgia"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = C_GOLD
        return title_box

    # 5. Helper to create cards
    def add_card(slide, left, top, width, height, border_color=C_GOLD, fill_color=C_BG_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # Blank layout is index 6 in standard templates
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: MUQOVA (PITCH DECK COVER)
    # ==========================================
    slide_1 = prs.slides.add_slide(blank_layout)
    background = slide_1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG_DARK
    
    # Golden border frame for cover
    frame = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7))
    frame.fill.background()
    frame.line.color.rgb = C_GOLD
    frame.line.width = Pt(2)
    
    # Large Floating Uzbek Mandala Plate on the Right
    try:
        slide_1.shapes.add_picture("uzbek_mandala.jpg", Inches(7.5), Inches(1.25), Inches(5.0), Inches(5.0))
    except Exception as e:
        pass
    
    # Logo on Cover
    try:
        slide_1.shapes.add_picture("logo.jpg", Inches(1.0), Inches(0.8), Inches(1.2), Inches(1.2))
    except Exception as e:
        pass
        
    # Title Box
    title_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(6.2), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "UZBEK-TRIP"
    p.font.name = "Georgia"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = C_GOLD
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Shaxsiy va Aqlli Marshrut Konstruktori"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(18)
    p2.font.color.rgb = C_WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "5 daqiqalik investitsion pitch deck taqdimoti: muammo, yechim, bozor, jamoa va erishilgan natijalar (traction)."
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(12)
    p3.font.color.rgb = C_SLATE_400
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(15)

    # Footer on Cover
    cov_foot = slide_1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(6.0), Inches(0.8))
    tf_f = cov_foot.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "INVESTOR PITCH DECK  |  2026"
    p_f.font.name = "Segoe UI"
    p_f.font.size = Pt(11)
    p_f.font.color.rgb = C_TURQUOISE
    p_f.font.bold = True

    # ==========================================
    # SLIDE 2: MUAMMO (THE PROBLEM)
    # ==========================================
    slide_2 = prs.slides.add_slide(blank_layout)
    apply_background(slide_2)
    add_slide_title(slide_2, "TURIZMDAGI ANIQ MUAMMOLAR")

    col_width = Inches(3.6)
    col_height = Inches(4.5)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(1.8)

    problems = [
        {
            "num": "01",
            "title": "Qotib qolgan turlar",
            "desc": "Mavjud turpaketlar juda qattiq rejalashtirilgan. Zamonaviy erkin sayyohlar tayyor shablonlardan charchagan va o'z marshrutlarini shaxsiy qiziqishlariga moslashni xohlashadi.",
            "sub": "Individual yondashuv yo'qligi"
        },
        {
            "num": "02",
            "title": "Noaniq va yashirin narxlar",
            "desc": "Gidlar va haydovchilar bozori tartibga solinmagan. Narxlar oldindan ma'lum emas va yashirin komissiyalar sababli sayyohlar ko'p pul yo'qotishadi.",
            "sub": "Moliyaviy shaffoflik yo'qligi"
        },
        {
            "num": "03",
            "title": "Sertifikatsiyalanmagan xizmatlar",
            "desc": "Professional tillarni biladigan gidlar va ishonchli haydovchilarni bitta platformada kafolatlangan holda band qilish va tekshirish imkoniyati cheklangan.",
            "sub": "Xavfsizlik va sifat nazorati sustligi"
        }
    ]

    for i, prob in enumerate(problems):
        left_pos = start_left + i * (col_width + gap)
        add_card(slide_2, left_pos, top_pos, col_width, col_height)
        
        tb = slide_2.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.2), col_width - Inches(0.4), col_height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_num = tf.paragraphs[0]
        p_num.text = prob["num"]
        p_num.font.name = "Georgia"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = C_TURQUOISE
        
        p_title = tf.add_paragraph()
        p_title.text = prob["title"]
        p_title.font.name = "Georgia"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = C_GOLD
        p_title.space_before = Pt(10)
        p_title.space_after = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = prob["desc"]
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = C_SLATE_400
        p_desc.space_after = Pt(10)
        
        p_sub = tf.add_paragraph()
        p_sub.text = prob["sub"]
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(11)
        p_sub.font.italic = True
        p_sub.font.color.rgb = C_SLATE_600

    # ==========================================
    # SLIDE 3: YECHIM: QANDAY ISHLAYDI (THE SOLUTION)
    # ==========================================
    slide_3 = prs.slides.add_slide(blank_layout)
    apply_background(slide_3)
    add_slide_title(slide_3, "YECHIM VA MAHSULOT QANDAY ISHLAYDI")

    left_tb = slide_3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_l = left_tb.text_frame
    tf_l.word_wrap = True
    
    pl_1 = tf_l.paragraphs[0]
    pl_1.text = "Biz sayyohlarga to'liq shaffoflik va boshqaruv erkinligini beramiz."
    pl_1.font.name = "Georgia"
    pl_1.font.size = Pt(24)
    pl_1.font.bold = True
    pl_1.font.color.rgb = C_WHITE
    pl_1.space_after = Pt(15)
    
    pl_2 = tf_l.add_paragraph()
    pl_2.text = "Sayyoh bir necha soniyalar ichida o'ziga mos transport turini, gid tili va boradigan maskanlarini tanlab, o'z yo'nalishini 'Lego' kabi teradi."
    pl_2.font.name = "Segoe UI"
    pl_2.font.size = Pt(13)
    pl_2.font.color.rgb = C_SLATE_400
    pl_2.space_after = Pt(15)

    pl_3 = tf_l.add_paragraph()
    pl_3.text = "Real vaqt rejimida GPS navigatsiya yo'llari chiziladi va jami narx hisoblanib chiqadi."
    pl_3.font.name = "Segoe UI"
    pl_3.font.size = Pt(12)
    pl_3.font.color.rgb = C_TURQUOISE

    right_left = Inches(5.8)
    card_w = Inches(6.7)
    card_h = Inches(1.3)
    card_gap = Inches(0.2)
    
    features = [
        {
            "title": "🗺 Interaktiv Marshrut Quruvchi",
            "desc": "Tarixiy diqqatga sazovor joylar va muqobil eko-turlarni birlashtirib, OSRM orqali yo'llar bo'yicha marshrut chizish."
        },
        {
            "title": "🚗 Haqiqiy GPS Kuzatuv & Simulyatsiya",
            "desc": "Sayyohning harakatlanishini online ko'rsatish, HUD panelda masofa, tezlik va yetib borish vaqtini (ETA) jonli yangilash."
        },
        {
            "title": "⚡ Avtomatlashtirilgan n8n Workflow",
            "desc": "Mijoz OTP kod orqali band qilganda haydovchi va gidning WhatsApp raqamiga to'liq ma'lumotlarni 1 soniyada yuborish."
        }
    ]

    for i, feat in enumerate(features):
        curr_top = Inches(1.8) + i * (card_h + card_gap)
        add_card(slide_3, right_left, curr_top, card_w, card_h)
        
        tb = slide_3.shapes.add_textbox(right_left + Inches(0.2), curr_top + Inches(0.1), card_w - Inches(0.4), card_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = feat["title"]
        p_t.font.name = "Georgia"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = C_GOLD
        
        p_d = tf.add_paragraph()
        p_d.text = feat["desc"]
        p_d.font.name = "Segoe UI"
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = C_SLATE_400
        p_d.space_before = Pt(4)

    # ==========================================
    # SLIDE 4: BOZOR HAJMI VA RAQOBAT (MARKET & COMPETITORS)
    # ==========================================
    slide_4 = prs.slides.add_slide(blank_layout)
    apply_background(slide_4)
    add_slide_title(slide_4, "BOZOR HAJMI VA ASOSIY RAQOBATCHILAR")

    box_w = Inches(3.6)
    box_h = Inches(2.2)
    top_pos_1 = Inches(1.8)
    top_pos_2 = Inches(4.3)
    left_1 = Inches(0.8)
    left_2 = Inches(4.8)
    left_3 = Inches(8.8)

    market_metrics = [
        {"val": "$1.2B+", "lbl": "TAM (Total Addressable Market)", "sub": "O'zbekistondagi umumiy turizm bozori va sayyohlar xarajatlari yillik hajmi."},
        {"val": "$250M", "lbl": "SAM (Serviceable Addressable Market)", "sub": "Mustaqil va individual ravishda sayohat qiluvchi xalqaro turistlar ulushi."},
        {"val": "$15M", "lbl": "SOM (Serviceable Obtainable Market)", "sub": "Dastlabki yillarda platforma orqali qamrab olinadigan tranzaksiyalar hajmi."}
    ]

    for i, met in enumerate(market_metrics):
        left_pos = left_1 + i * (box_w + Inches(0.4))
        add_card(slide_4, left_pos, top_pos_1, box_w, box_h, border_color=C_TURQUOISE)
        
        tb = slide_4.shapes.add_textbox(left_pos + Inches(0.2), top_pos_1 + Inches(0.15), box_w - Inches(0.4), box_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_val = tf.paragraphs[0]
        p_val.text = met["val"]
        p_val.font.name = "Georgia"
        p_val.font.size = Pt(40)
        p_val.font.bold = True
        p_val.font.color.rgb = C_GOLD
        p_val.alignment = PP_ALIGN.CENTER
        
        p_lbl = tf.add_paragraph()
        p_lbl.text = met["lbl"]
        p_lbl.font.name = "Segoe UI"
        p_lbl.font.size = Pt(12)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = C_WHITE
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.space_before = Pt(5)

        p_sub = tf.add_paragraph()
        p_sub.text = met["sub"]
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = C_SLATE_400
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_before = Pt(5)

    panel_w = Inches(11.733)
    panel_h = Inches(2.0)
    add_card(slide_4, left_1, top_pos_2, panel_w, panel_h)
    
    tb_p = slide_4.shapes.add_textbox(left_1 + Inches(0.3), top_pos_2 + Inches(0.2), panel_w - Inches(0.6), panel_h - Inches(0.4))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    
    pp_1 = tf_p.paragraphs[0]
    pp_1.text = "⚔ Raqobat Muhiti va Bizning Ustunligimiz (USP)"
    pp_1.font.name = "Georgia"
    pp_1.font.size = Pt(17)
    pp_1.font.bold = True
    pp_1.font.color.rgb = C_GOLD
    pp_1.space_after = Pt(6)
    
    pp_2 = tf_p.add_paragraph()
    pp_2.text = "Raqobatchilarimiz (an'anaviy turoperatorlar va ko'cha taksi xizmatlari) qimmat va rejalari qat'iy belgilangan turlarni taklif qiladilar. Bizning UZBEK-TRIP platformamiz esa sayyohlarga o'z turlarini 100% mustaqil tuzish, shaffof hisoblangan narxni ko'rish, onlayn GPS navigatsiya va n8n yordamida tezkor WhatsApp/SMS tasdiqnomaga ega bo'lish ustunligini beradi."
    pp_2.font.name = "Segoe UI"
    pp_2.font.size = Pt(11.5)
    pp_2.font.color.rgb = C_SLATE_400

    # ==========================================
    # SLIDE 5: JAMOA VA NATIJALAR (TEAM & TRACTION)
    # ==========================================
    slide_5 = prs.slides.add_slide(blank_layout)
    apply_background(slide_5)
    add_slide_title(slide_5, "JAMOA VA ERISHILGAN NATIJALAR (TRACTION)")

    col_w = Inches(5.6)
    col_h = Inches(4.5)
    gap = Inches(0.5)
    start_left = Inches(0.8)
    top_pos = Inches(1.8)

    # Left Card: Team
    add_card(slide_5, start_left, top_pos, col_w, col_h)
    tb_t = slide_5.shapes.add_textbox(start_left + Inches(0.3), top_pos + Inches(0.3), col_w - Inches(0.6), col_h - Inches(0.6))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True

    pt_t = tf_t.paragraphs[0]
    pt_t.text = "👥 Professional Jamoa"
    pt_t.font.name = "Georgia"
    pt_t.font.size = Pt(20)
    pt_t.font.bold = True
    pt_t.font.color.rgb = C_GOLD
    pt_t.space_after = Pt(15)

    team_members = [
        ("Muhriddin Bo'riboyev", "Founder & Lead Developer. Next.js, Supabase, OSRM routing va n8n avtomatlashtirish muhandisi."),
        ("Operatsion Guruh", "Mehmonxonalar integratsiyasi, gidlar va haydovchilar bilan hamkorlik tarmoqlarini muvofiqlashtiruvchi operatsion menejerlar."),
        ("Lokal Hamkorlar", "Samarqand, Buxoro va Toshkent shaharlaridagi litsenziyalangan gidlar va sertifikatlangan transport uyushmalari.")
    ]
    for name, role in team_members:
        p = tf_t.add_paragraph()
        p.text = f"• {name} - "
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
        r = p.add_run()
        r.text = role
        r.font.bold = False
        r.font.color.rgb = C_SLATE_400
        p.space_after = Pt(12)

    # Right Card: Traction
    right_left = start_left + col_w + gap
    add_card(slide_5, right_left, top_pos, col_w, col_h)
    tb_tr = slide_5.shapes.add_textbox(right_left + Inches(0.3), top_pos + Inches(0.3), col_w - Inches(0.6), col_h - Inches(0.6))
    tf_tr = tb_tr.text_frame
    tf_tr.word_wrap = True

    ptr_t = tf_tr.paragraphs[0]
    ptr_t.text = "📈 Erishilgan Natijalar (Traction)"
    ptr_t.font.name = "Georgia"
    ptr_t.font.size = Pt(20)
    ptr_t.font.bold = True
    ptr_t.font.color.rgb = C_GOLD
    ptr_t.space_after = Pt(15)

    tractions = [
        ("Ishlovchi MVP", "Veb platformamiz Vercel-ga joylashtirildi va to'liq foydalanishga tayyor: uzbek-trip.vercel.app"),
        ("OSRM va GPS Navigatsiyasi", "Xaritada yo'l tarmog'i bo'yicha marshrut hisoblash, live navigatsiya va virtual sayohat simulyatori joriy qilindi."),
        ("n8n Avtomatizatsiyasi", "Sayohat buyurtmalarini operatsion tasdiqlash va haydovchi/gidlarga WhatsApp orqali yetkazish moduli sinovdan o'tdi."),
        ("B2B Hamkorlik Kelishuvlari", "Toshkent va Samarqand shahridagi dastlabki mehmonxonalar va gidlar bilan integratsiya bo'yicha kelishuvlar.")
    ]
    for lbl, desc in tractions:
        p = tf_tr.add_paragraph()
        p.text = f"✔ {lbl}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_TURQUOISE
        
        r = p.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = C_SLATE_400
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 6: BIZNES MODEL VA KELAJAK REJALAR (BUSINESS MODEL)
    # ==========================================
    slide_6 = prs.slides.add_slide(blank_layout)
    apply_background(slide_6)
    add_slide_title(slide_6, "BIZNES MODEL VA STRATEGIK KELAJAK")

    col_w = Inches(3.6)
    col_h = Inches(4.5)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(1.8)

    plans = [
        {
            "title": "Komissiya & Daromad",
            "val": "15% platforma ulushi",
            "desc": "Platforma orqali band qilingan har bir transport va gid xizmatlaridan 15% tranzaksion komissiya to'lovi olinadi. O'rtacha buyurtma qiymati $120.",
            "points": [
                "15% gid va transport komissiyasi",
                "Mehmonxona va restoran keshbeklari"
            ]
        },
        {
            "title": "UzRailways API Integratsiya",
            "desc": "Kelgusi oylarda sayyohlar uchun platformaning o'zida Afrosiyob tezyurar poyezdi va avia-chiptalarni birdan xarid qilish tizimi qo'shiladi.",
            "val": "Kelasi reja: chiptalar",
            "points": [
                "Afrosiyob poyezd chiptalari",
                "Lokal aviareys integratsiyasi"
            ]
        },
        {
            "title": "Hududiy Kengayish",
            "desc": "Platformani Samarqanddan so'ng Toshkent, Buxoro, Xiva va Qoraqalpog'iston viloyatlariga to'liq moslashtirish va sayohatlarni viloyatlararo ulash.",
            "val": "Kelasi reja: kengayish",
            "points": [
                "Buxoro, Xiva va Toshkent",
                "Viloyatlararo turlar"
            ]
        }
    ]

    for i, plan in enumerate(plans):
        left_pos = start_left + i * (col_w + gap)
        add_card(slide_6, left_pos, top_pos, col_w, col_h)
        
        tb = slide_6.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.2), col_w - Inches(0.4), col_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = plan["title"]
        p_t.font.name = "Georgia"
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = C_GOLD
        p_t.space_after = Pt(10)
        
        p_val = tf.add_paragraph()
        p_val.text = plan["val"]
        p_val.font.name = "Georgia"
        p_val.font.size = Pt(20)
        p_val.font.bold = True
        p_val.font.color.rgb = C_TURQUOISE
        p_val.space_after = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = plan["desc"]
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = C_WHITE
        p_desc.space_after = Pt(15)
        
        for pt in plan["points"]:
            p_pt = tf.add_paragraph()
            p_pt.text = "• " + pt
            p_pt.font.name = "Segoe UI"
            p_pt.font.size = Pt(11)
            p_pt.font.color.rgb = C_SLATE_400
            p_pt.space_after = Pt(5)

    # ==========================================
    # SLIDE 7: SAVOLLAR VA BOG'LANISH (Q&A)
    # ==========================================
    slide_7 = prs.slides.add_slide(blank_layout)
    apply_background(slide_7)
    
    # Large Decorative Golden Frame in Center
    frame = slide_7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(1.5), Inches(9.333), Inches(4.5))
    frame.fill.solid()
    frame.fill.fore_color.rgb = C_BG_CARD
    frame.line.color.rgb = C_GOLD
    frame.line.width = Pt(2)
    
    tb = slide_7.shapes.add_textbox(Inches(2.5), Inches(2.0), Inches(8.333), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "KATTA RAHMAT! SAVOLLAR?"
    p1.font.name = "Georgia"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = C_GOLD
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(15)
    
    p2 = tf.add_paragraph()
    p2.text = "O'zbekistonda shaxsiy turizm kelajagini birgalikda barpo etamiz."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.color.rgb = C_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(25)
    
    p3 = tf.add_paragraph()
    p3.text = "📞 Telefon: +998 (94) 019-64-20\n📧 Email: mikemukhriddin@gmail.com\n🌐 Veb-sayt: uzbek-trip.vercel.app"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_SLATE_400
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(15)

    p4 = tf.add_paragraph()
    p4.text = "5 daqiqada qisqa va lo'nda javob berishga tayyormiz!"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_TURQUOISE
    p4.font.bold = True
    p4.font.italic = True
    p4.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_filename = "Samarqand_CrafTour_Investor_Pitch_Deck.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'!")

    # Convert to PDF if possible
    try:
        import os
        import win32com.client
        
        print("Converting presentation to PDF...")
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        pptx_abs = os.path.abspath(output_filename)
        pdf_filename = "Samarqand_CrafTour_Investor_Pitch_Deck.pdf"
        pdf_abs = os.path.abspath(pdf_filename)
        
        deck = powerpoint.Presentations.Open(pptx_abs, WithWindow=False)
        deck.SaveAs(pdf_abs, 32)
        deck.Close()
        powerpoint.Quit()
        print(f"Presentation successfully converted and saved as '{pdf_filename}'!")
    except Exception as e:
        print(f"Notice: Could not automatically convert PPTX to PDF via win32com: {e}")

if __name__ == "__main__":
    create_deck()
